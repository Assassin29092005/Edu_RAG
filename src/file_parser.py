"""
file_parser.py — Extracts text from PDF, DOCX, and PPTX files.
Returns a list of document dicts: { text, source, page/slide }
"""

import os
from PIL import Image
import fitz  # PyMuPDF
import pytesseract
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from unstructured.partition.docx import partition_docx
from unstructured.partition.pptx import partition_pptx
from unstructured.chunking.title import chunk_by_title

from src.vision_utils import summarize_image

# Configure Tesseract Path (for Windows)
try:
    pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
except Exception:
    pass

def parse_pdf(file_path: str) -> list[dict]:
    """Extract text page-by-page from a PDF file. Falls back to OCR for scanned pages."""
    docs = []
    pdf = fitz.open(file_path)
    filename = os.path.basename(file_path)

    for page_num in range(len(pdf)):
        page = pdf[page_num]
        text = page.get_text().strip()
        
        # If less than 50 chars natively, assume it's a scanned page / image
        if len(text) < 50:
            pix = page.get_pixmap(dpi=300) # high res for OCR
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            
            try:
                ocr_text = pytesseract.image_to_string(img)
                text = f"{text}\n{ocr_text}".strip()
            except Exception as e:
                print(f"OCR failed on {filename} page {page_num+1}: {e}")
        
        # Extract images from PDF and summarize them using Multimodal Vision
        try:
            images_dir = os.path.join(os.path.dirname(__file__), "..", "data", "images")
            os.makedirs(images_dir, exist_ok=True)
            safe_filename = "".join([c if c.isalnum() else "_" for c in filename])

            for img_index, img in enumerate(page.get_images(full=True)):
                xref = img[0]
                base_image = pdf.extract_image(xref)
                image_bytes = base_image["image"]
                
                # Save image to disk
                image_ext = base_image["ext"]
                image_filename = f"{safe_filename}_page{page_num+1}_img{img_index}.{image_ext}"
                image_path = os.path.join(images_dir, image_filename)
                with open(image_path, "wb") as f:
                    f.write(image_bytes)

                summary = summarize_image(image_bytes)
                if summary:
                    text += summary
        except Exception as e:
             print(f"PDF Image extraction failed: {e}")
                
        if text:
            docs.append({
                "text": text,
                "source": filename,
                "page": page_num + 1,
                "type": "pdf"
            })

    pdf.close()
    return docs

def parse_docx(file_path: str) -> list[dict]:
    """Extract text semantically from a DOCX file using Unstructured."""
    filename = os.path.basename(file_path)
    docs = []
    
    try:
        elements = partition_docx(filename=file_path)
        chunks = chunk_by_title(elements) # Group into semantic parent documents
        
        for i, chunk in enumerate(chunks, 1):
            if chunk.text.strip():
                docs.append({
                    "text": chunk.text,
                    "source": filename,
                    "page": i, # treat chunk ID as 'page' conceptually for reference
                    "type": "docx"
                })
    except Exception as e:
        print(f"Failed to parse docx {filename}: {e}")

    return docs

def parse_pptx(file_path: str) -> list[dict]:
    """Extract text semantically from a PPTX file and summarize images."""
    filename = os.path.basename(file_path)
    docs = []

    # 1. Extract image summaries per slide using python-pptx and OCR
    slide_summaries = {}
    try:
        prs = Presentation(file_path)
        
        # Prepare directory for saving extracted images
        images_dir = os.path.join(os.path.dirname(__file__), "..", "data", "images")
        os.makedirs(images_dir, exist_ok=True)
        safe_filename = "".join([c if c.isalnum() else "_" for c in filename])

        for i, slide in enumerate(prs.slides, 1):
            slide_summaries[i] = ""
            img_index = 0
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_bytes = shape.image.blob
                    
                    # Save image to disk
                    image_ext = shape.image.ext
                    image_filename = f"{safe_filename}_page{i}_img{img_index}.{image_ext}"
                    image_path = os.path.join(images_dir, image_filename)
                    with open(image_path, "wb") as f:
                        f.write(image_bytes)
                    img_index += 1

                    # Extract text physically baked into the image using OCR
                    try:
                        import io
                        img = Image.open(io.BytesIO(image_bytes))
                        ocr_text = pytesseract.image_to_string(img).strip()
                        if ocr_text:
                            slide_summaries[i] += f"\n[Diagram Text OCR: {ocr_text}]\n"
                    except Exception as e:
                        print(f"OCR Failed on image: {e}")

                    # Generate semantic summary using Vision Model
                    summary = summarize_image(image_bytes)
                    if summary:
                        slide_summaries[i] += summary
    except Exception as e:
        print(f"Error parsing PPTX images {filename}: {e}")

    # 2. Extract semantic text using Unstructured and group strictly by slide/page
    try:
        from unstructured.partition.pptx import partition_pptx
        
        elements = partition_pptx(filename=file_path)
        
        # Group text by page number
        slide_texts = {}
        for el in elements:
            page_num = el.metadata.page_number
            if not page_num:
                page_num = 1 # Fallback
                
            if page_num not in slide_texts:
                slide_texts[page_num] = []
            
            if el.text.strip():
                slide_texts[page_num].append(el.text.strip())
                
        # Combine extracted text with generated image summaries for each slide
        all_pages = set(slide_texts.keys()).union(set(slide_summaries.keys()))
        
        for page_num in sorted(list(all_pages)):
            combined_text = ""
            
            if page_num in slide_texts:
                combined_text += "\n".join(slide_texts[page_num]) + "\n"
                
            if page_num in slide_summaries and slide_summaries[page_num]:
                combined_text += "\n" + slide_summaries[page_num] + "\n"
                
            combined_text = combined_text.strip()
            
            if combined_text:
                docs.append({
                    "text": combined_text,
                    "source": filename,
                    "page": page_num,
                    "type": "pptx"
                })
    except Exception as e:
        print(f"Failed to parse pptx {filename}: {e}")

    return docs

def parse_file(file_path: str) -> list[dict]:
    """Parse a file based on its extension. Returns list of document dicts."""
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return parse_pdf(file_path)
    elif ext == ".docx":
        return parse_docx(file_path)
    elif ext in (".ppt", ".pptx"):
        return parse_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file format: {ext}")
