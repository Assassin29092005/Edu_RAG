from pptx import Presentation
prs = Presentation('data/uploads/VIRTUAL MACHINES PROVISIONING_UNIT2_PART1.pptx')

for i, slide in enumerate(prs.slides, 1):
    if 5 <= i <= 8:
        print(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(f"TEXT: {repr(shape.text[:50])}")
            print(f"SHAPE: {shape.shape_type}")
